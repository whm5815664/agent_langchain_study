import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

output_dir = os.environ.get('OUTPUT_DIR', r'E:\code\jupyter\智能体学习\langchain\outputs')
output_path = os.path.join(output_dir, 'skills.pptx')

print('Output path:', output_path)

prs = Presentation()
primary_color = RGBColor(2, 128, 144)
dark_bg = RGBColor(30, 39, 97)

skills = [
    {'name': 'docx', 'description': 'Word 文档创建、编辑和分析', 'details': ['创建新文档', '编辑现有文档', '支持修订跟踪', '文本提取', '支持表格图片']},
    {'name': 'pdf', 'description': 'PDF 文件处理', 'details': ['合并/拆分 PDF', '文本表格提取', '创建新 PDF', 'OCR 扫描 PDF', '添加水印']},
    {'name': 'pptx', 'description': 'PowerPoint 演示文稿处理', 'details': ['创建和编辑', '使用模板', '文本提取', '设计指南', '视觉 QA']},
    {'name': 'skill-creator', 'description': '创建有效技能的指南', 'details': ['技能结构', '设计模式', '创建流程', '资源管理', '打包验证']},
    {'name': 'xlsx', 'description': '电子表格文件处理', 'details': ['数据分析', '创建/编辑', '颜色编码', '公式计算', '数据格式']}
]

# Slide 1: Title
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'Skills 技能库'

# Slide 2: Overview
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '技能库概览'

# Slides 3-7: Individual Skills
for skill in skills:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = skill['name'].upper()

# Slide 8: Summary
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '总结'

prs.save(output_path)
print('Saved to:', output_path)
print('File size:', os.path.getsize(output_path), 'bytes')

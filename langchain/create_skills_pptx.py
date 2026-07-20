import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Get output directory
output_dir = os.environ.get('OUTPUT_DIR', 'outputs')
output_path = os.path.join(output_dir, 'skills.pptx')

# Create presentation
prs = Presentation()

# Define color palette (Teal Trust theme)
primary_color = RgbColor(2, 128, 144)    # Teal
secondary_color = RgbColor(0, 168, 150)   # Seafoam
accent_color = RgbColor(2, 195, 154)      # Mint
dark_bg = RgbColor(30, 39, 97)            # Navy
white = RgbColor(255, 255, 255)

# Skills data
skills = [
    {
        'name': 'docx',
        'description': 'Word 文档创建、编辑和分析',
        'details': [
            '创建新文档（使用 docx-js）',
            '编辑现有文档（解包→编辑 XML→打包）',
            '支持修订跟踪和批注',
            '文本提取和格式保留',
            '支持表格、图片、目录等'
        ]
    },
    {
        'name': 'pdf',
        'description': 'PDF 文件处理',
        'details': [
            '合并/拆分 PDF 文件',
            '文本和表格提取',
            '创建新 PDF（使用 reportlab）',
            'OCR 扫描的 PDF',
            '添加水印和密码保护',
            '填写 PDF 表单'
        ]
    },
    {
        'name': 'pptx',
        'description': 'PowerPoint 演示文稿处理',
        'details': [
            '创建和编辑演示文稿',
            '使用模板或从头创建',
            '文本提取和分析',
            '设计指南和配色方案',
            '视觉 QA 和验证'
        ]
    },
    {
        'name': 'skill-creator',
        'description': '创建有效技能的指南',
        'details': [
            '技能结构和组织',
            '渐进式披露设计模式',
            '技能创建流程（6 步）',
            '可重用资源管理',
            '技能打包和验证'
        ]
    },
    {
        'name': 'xlsx',
        'description': '电子表格文件处理',
        'details': [
            '使用 pandas 进行数据分析',
            '使用 openpyxl 创建/编辑',
            '财务模型颜色编码标准',
            '公式 recalculating（LibreOffice）',
            '数据格式和验证'
        ]
    }
]

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Skills 技能库"
subtitle.text = "专业任务处理工具集\n文档 • PDF • 演示文稿 • 电子表格 • 技能创建"

# Style title slide
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = dark_bg
subtitle.text_frame.paragraphs[0].font.size = Pt(20)
subtitle.text_frame.paragraphs[0].font.color.rgb = secondary_color

# Slide 2: Overview
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "技能库概览"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = dark_bg

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

overview_text = [
    "5 个专业技能模块",
    "覆盖常见办公文档格式",
    "提供创建、编辑、分析功能",
    "包含设计指南和最佳实践",
    "支持自动化工作流"
]

for i, item in enumerate(overview_text):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(18)
    p.level = 0

# Slide 3-7: Individual Skills
for skill in skills:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = skill['name'].upper()
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Add description
    p = tf.paragraphs[0]
    p.text = skill['description']
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = dark_bg
    
    # Add details
    for i, detail in enumerate(skill['details']):
        p = tf.add_paragraph()
        p.text = detail
        p.font.size = Pt(16)
        p.level = 0
        p.space_before = Pt(12)

# Slide 8: Summary
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "总结"
title.text_frame.paragraphs[0].font.size = Pt(36)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = dark_bg

content = slide.placeholders[1]
tf = content.text_frame
tf.clear()

summary_points = [
    "完整的文档处理工具集",
    "支持主流办公格式",
    "遵循行业最佳实践",
    "可扩展的技能架构",
    "自动化与人工审核结合"
]

for i, point in enumerate(summary_points):
    p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
    p.text = point
    p.font.size = Pt(20)
    p.level = 0

# Save presentation
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
